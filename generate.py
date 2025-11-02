import os
import random
import json
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import requests  # For downloading images
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

def generate_presentation(topic_category, num_slides):
    # Create outputs directory if it doesn't exist
    outputs_dir = "outputs"
    os.makedirs(outputs_dir, exist_ok=True)
    
    # Set up OpenAI client
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        api_key = "OR_HARDCODED_OPENAI_API_KEY"
    client = OpenAI(api_key=api_key)

    # Prompt for generating a random topic within the category
    topic_prompt = f"""
    Generate a single random niche topic within the category '{topic_category}'.
    The topic should be suitable for an educational presentation aimed at a smart college student.
    Make it interesting and provide a balance with more technical depth.
    Output only the topic title as a string, nothing else.
    """

    # Call OpenAI API for topic
    topic_response = client.chat.completions.create(
        model="gpt-4o",  # Use a cost-effective model; can change to gpt-4o for better quality
        messages=[
            {"role": "system", "content": "You are a helpful assistant for creating educational topics."},
            {"role": "user", "content": topic_prompt}
        ]
    )

    topic = topic_response.choices[0].message.content.strip()
    print(f"Generating presentation on: {topic} with {num_slides} slides.")

    # Prompt for generating slide content
    slides_prompt = f"""
    Create a {num_slides}-slide presentation outline on the topic '{topic}'.
    The audience is a smart college student. 
    Aim to explain the topic with more technical depth and detailed explanations.
    Include a title slide as slide 1, content slides with expanded details, and a conclusion slide if it fits.
    Output ONLY the following JSON format, nothing else - no explanations or additional text:
    [
        {{"slide": 1, "title": "Title Here", "content": ["bullet1", "bullet2"]}},
        ...
    ]
    Ensure content is educational, engaging, accurate, and detailed.
    """

    # Call OpenAI API for slides
    slides_response = client.chat.completions.create(
        model="gpt-4o-mini",  # Use a cost-effective model; can change to gpt-4o for better quality
        messages=[
            {"role": "system", "content": "You are a helpful assistant for creating educational content."},
            {"role": "user", "content": slides_prompt}
        ]
    )

    # Parse the JSON response
    try:
        slides_data = json.loads(slides_response.choices[0].message.content)
    except json.JSONDecodeError as e:
        print(f"Error parsing OpenAI response: {e}")
        print("Raw response:")
        print(slides_response.choices[0].message.content)
        return

    # Create PowerPoint with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add slides based on data
    for slide_info in slides_data:
        if slide_info["slide"] == 1:
            # Title slide layout (usually index 0)
            slide_layout = prs.slide_layouts[0]
        else:
            # Content slide layout (title and content, usually index 1)
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        title_p = title_shape.text_frame.paragraphs[0]
        title_p.font.size = Pt(44)
        title_p.font.bold = True

        if slide_info["slide"] == 1:
            title_p.alignment = PP_ALIGN.CENTER
            # Adjust title position and size to leave space for subtitle
            title_shape.top = Inches(1)
            title_shape.height = Inches(2)
            title_shape.width = Inches(10)
            title_shape.left = Inches(1.666)  # Center in 13.333 width
            title_shape.text_frame.word_wrap = True
        else:
            title_p.alignment = PP_ALIGN.LEFT
            title_shape.left = Inches(1)
            title_shape.width = Inches(11.333)  # Full width to prevent vertical stacking
            title_shape.height = Inches(1.5)  # Adequate height for wrapping
            title_shape.top = Inches(0.5)  # Position appropriately
            title_shape.text_frame.word_wrap = True

        if slide_info["slide"] != 1:  # Skip image for title slide
            # Generate image
            image_prompt = f"A nice, simple, illustrative image for a presentation slide titled '{slide_info['title']}' in the context of {topic_category.lower()}, educational, photo realistic (where possible)."
            image_response = client.images.generate(
                model="dall-e-3",
                prompt=image_prompt,
                size="1024x1024",
                quality="standard",
                n=1
            )
            image_url = image_response.data[0].url
            # Download image
            img_data = requests.get(image_url).content
            img_filename = os.path.join(outputs_dir, f"slide_{slide_info['slide']}_image.png")
            with open(img_filename, 'wb') as handler:
                handler.write(img_data)

            # Insert image on right side, right-aligned with margin
            img_width = Inches(4.2)  # 70% of previous 6 inches
            img_left = prs.slide_width - img_width - Inches(0.5)  # Right margin 0.5 inch
            img_top = Inches(2)
            slide.shapes.add_picture(img_filename, img_left, img_top, width=img_width)  # Maintain ratio

            # Clean up
            os.remove(img_filename)

        if "content" in slide_info and slide_info["content"]:
            if slide_info["slide"] == 1 and len(slide.placeholders) > 1:
                # Subtitle for title slide
                subtitle_shape = slide.placeholders[1]
                tf = subtitle_shape.text_frame
                tf.clear()  # Clear default text

                # Adjust subtitle position to below title without overlap
                subtitle_shape.top = Inches(3.5)
                subtitle_shape.height = Inches(3)
                subtitle_shape.width = Inches(10)
                subtitle_shape.left = Inches(1.666)  # Center
            else:
                # Content placeholder for other slides (usually index 1)
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()  # Clear default text

            tf.word_wrap = True
            for i, bullet in enumerate(slide_info["content"]):
                if i > 0 or not tf.text:  # Add paragraphs after first
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                p.text = bullet
                p.font.size = Pt(24)
                p.font.bold = False
                p.alignment = PP_ALIGN.LEFT if slide_info["slide"] != 1 else PP_ALIGN.CENTER
                p.level = i  # Hierarchical bullets

            if slide_info["slide"] != 1:
                # Resize text to fill left space up to image
                text_width = img_left - Inches(1) - Inches(0.5)
                content_shape.left = Inches(1)
                content_shape.top = Inches(2)
                content_shape.width = text_width
                content_shape.height = Inches(4.5)

    # Save the presentation as PPTX
    pptx_filename = os.path.join(outputs_dir, f"{topic.replace(' ', '_')}_{num_slides}_slides.pptx")
    prs.save(pptx_filename)
    print(f"PowerPoint saved as: {pptx_filename}")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Generate a random PowerPoint presentation using OpenAI and python-pptx.")
    parser.add_argument("--topic_category", default="Geography", help="Category for the random topic. Defaults to 'Geography' if not provided.")
    parser.add_argument("--num_slides", choices=[5, 10, 15], type=int, required=True, help="Number of slides in the presentation.")
    args = parser.parse_args()

    generate_presentation(args.topic_category, args.num_slides)